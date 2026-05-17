from __future__ import annotations

from io import BytesIO
from pathlib import PurePosixPath


def _extension(filename: str) -> str:
    return PurePosixPath(filename).suffix.lower().lstrip(".")


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


def extract_text_from_bytes(
    *,
    filename: str,
    content: bytes,
    content_type: str | None = None,
) -> str:
    """Extract searchable text from common uploaded course/admin document files."""
    ext = _extension(filename)
    mime = content_type or ""

    if ext in {"txt", "md", "csv"} or mime.startswith("text/"):
        return _decode_text(content)

    if ext == "pdf" or mime == "application/pdf":
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise RuntimeError("pypdf is required to extract PDF text") from exc

        reader = PdfReader(BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages).strip()

    if ext == "docx":
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("python-docx is required to extract DOCX text") from exc

        document = Document(BytesIO(content))
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()

    if ext in {"ppt", "pptx"}:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("python-pptx is required to extract PPT/PPTX text") from exc

        presentation = Presentation(BytesIO(content))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    lines.append(text)
        return "\n".join(lines).strip()

    return _decode_text(content).strip()
